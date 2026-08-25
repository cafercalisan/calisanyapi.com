from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "calisan-yapi-rakip-arastirmasi-2026-08.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

INK = RGBColor(20, 31, 34)
PAPER = RGBColor(239, 235, 224)
TEAL = RGBColor(27, 171, 160)
GOLD = RGBColor(216, 160, 45)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(91, 105, 106)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, x, y, w, h, size=20, color=INK, bold=False,
            font="Aptos", align=PP_ALIGN.LEFT, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, kicker, heading, note=None):
    textbox(slide, kicker.upper(), .7, .42, 5, .35, 9, TEAL, True)
    textbox(slide, heading, .7, .86, 11.8, 1.05, 31, INK, True, "Aptos Display")
    if note:
        textbox(slide, note, .72, 1.82, 11.5, .55, 11, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.7), Inches(2.32), Inches(11.9), Inches(.025))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(198, 195, 184); line.line.fill.background()


def footer(slide, number, source=None):
    textbox(slide, "UMA YAPI · RAKİP ARAŞTIRMASI", .7, 7.12, 5, .2, 7, MUTED, True)
    textbox(slide, f"{number:02d}", 12.05, 7.05, .55, .24, 9, TEAL, True, align=PP_ALIGN.RIGHT)
    if source:
        textbox(slide, source, 6.2, 7.05, 5.5, .27, 6.5, MUTED, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(10); p.level = 0
        p.text = "•  " + p.text
    return box


def card(slide, x, y, w, h, heading, body, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(205, 202, 191)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.06), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    textbox(slide, heading, x+.22, y+.2, w-.42, .38, 14, INK, True)
    textbox(slide, body, x+.22, y+.7, w-.42, h-.86, 10, MUTED)


def image_cover(slide, path, x, y, w, h):
    from PIL import Image
    p = Path(path)
    with Image.open(p) as im:
        iw, ih = im.size
    target = w / h
    ratio = iw / ih
    if ratio > target:
        crop_w = ih * target
        left = int((iw - crop_w) / 2)
        crop = (left, 0, int(left + crop_w), ih)
    else:
        crop_h = iw / target
        top = int((ih - crop_h) / 2)
        crop = (0, top, iw, int(top + crop_h))
    with Image.open(p) as im:
        tmp = Path("/tmp") / f"crop-{p.stem}.png"
        im.crop(crop).save(tmp)
    slide.shapes.add_picture(str(tmp), Inches(x), Inches(y), Inches(w), Inches(h))


# 01 Cover
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, INK)
textbox(s, "PAZAR · WEB · SOSYAL MEDYA", .75, .62, 5.5, .4, 10, TEAL, True)
textbox(s, "Rakip görünürlüğünden\nUma Yapı’nın fırsat alanına.", .75, 1.25, 8.5, 2.1, 34, WHITE, True, "Aptos Display")
textbox(s, "İstanbul yapı, cam sistemleri ve sineklik pazarı\nKaynaklı masa başı araştırması · 10 Ağustos 2026", .78, 3.65, 6.8, .8, 14, RGBColor(190,200,199))
image_cover(s, ROOT/"public/brand/calisan-yapi-hero.png", 8.35, 0, 4.98, 7.5)
footer(s, 1)

# 02 Scope
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Araştırma çerçevesi","Neyi taradık, neyi doğruladık?")
card(s,.7,2.65,3.72,3.55,"WEB / SEO","Exa ile hizmet bazlı arama; Jina ile resmi sitelerin navigasyon, ürün, güven, teklif ve galeri yapılarının okunması.")
card(s,4.8,2.65,3.72,3.55,"SOSYAL / VİDEO","Resmi sosyal profil bağlantıları ve YouTube arama sonuçları. Instagram akışı OpenCLI eklentisi bağlı olmadığı için doğrudan okunamadı.",GOLD)
card(s,8.9,2.65,3.72,3.55,"GÖRSEL / ETİK","Rakip görselleri yalnız analiz referansıdır. Site veya reklamda yeniden kullanmayın; aynı kullanım senaryolarını kendi projelerinizle çekin.")
footer(s,2,"Agent Reach: Exa + Jina Reader + yt-dlp")

# 03 shortlist
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Rakip haritası","Takip edilmesi gereken 9 marka")
rows = [
    ("Albert Genau","Ulusal sistem evi","Teknik güven + ürün seçici"),
    ("Royal Glass","Üretici / bayi","Ürün ailesi + süreç"),
    ("Usta Cam Balkon","İstanbul uygulamacısı","Hız, taksit, garanti"),
    ("İpek Yapı","İstanbul uygulamacısı","Müşteri itirazı + referans"),
    ("Plisseri","Sineklik üreticisi","AR-GE + ihracat"),
    ("EM-SA","Sineklik uygulamacısı","Model çeşitliliği + yerel SEO"),
    ("Ricco Yapı","Premium dış mekân","Mimari yaşam dili"),
    ("Winter Garden","Premium kategori markası","İlçe/use-case segmentasyonu"),
    ("Ehil Yapı","Geniş hizmet uygulamacısı","Gerçek proje hikâyeleri"),
]
for i,(name,typ,edge) in enumerate(rows):
    col=i%3; row=i//3; x=.7+col*4.05; y=2.55+row*1.27
    card(s,x,y,3.7,1.03,name,f"{typ} · {edge}",TEAL if col!=1 else GOLD)
footer(s,3)

# 04 matrix
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Benchmark","Kim hangi alanda güçlü?")
headers=["Marka","Teknik içerik","Gerçek proje","Dönüşüm","Yerel SEO","Sosyal video"]
data=[
 ["Albert Genau","5","4","4","4","5"],["Royal Glass","4","3","3","2","3"],
 ["Usta Cam Balkon","3","3","5","3","2"],["İpek Yapı","2","4","4","2","2"],
 ["Plisseri","4","3","2","2","3"],["EM-SA","3","4","4","5","2"],
 ["Ricco / Winter","4","5","3","4","3"],["Ehil Yapı","3","5","3","3","1"],
]
x0=.72; y0=2.55; widths=[2.55,1.72,1.72,1.72,1.72,1.72]
for c,h in enumerate(headers):
    x=x0+sum(widths[:c]); textbox(s,h,x,y0,widths[c],.43,10,WHITE,True,align=PP_ALIGN.CENTER)
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y0),Inches(widths[c]),Inches(.48)); sh.fill.solid(); sh.fill.fore_color.rgb=INK; sh.line.color.rgb=PAPER
    textbox(s,h,x,y0+.08,widths[c],.3,9,WHITE,True,align=PP_ALIGN.CENTER)
for r,row in enumerate(data):
    for c,val in enumerate(row):
        x=x0+sum(widths[:c]); y=y0+.48+r*.48
        sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(widths[c]),Inches(.48)); sh.fill.solid(); sh.fill.fore_color.rgb=WHITE if r%2==0 else RGBColor(232,229,218); sh.line.color.rgb=RGBColor(210,207,196)
        color=INK if c==0 else (TEAL if int(val)>=4 else GOLD if int(val)==3 else MUTED)
        textbox(s,val,x,y+.08,widths[c],.28,9,color,c>0,align=PP_ALIGN.CENTER)
textbox(s,"Puanlar göreli masa başı değerlendirmedir (1 zayıf · 5 güçlü); trafik veya gelir tahmini değildir.",.75,6.65,10,.25,8,MUTED)
footer(s,4)

# 05 visual references
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Görsel benchmark","Gerçek uygulama sunumunda öne çıkan yaklaşım")
pics=[("Royal Glass · giyotin","/tmp/royal.jpg","royalglass.com.tr"),("İpek Yapı · kış bahçesi","/tmp/ipek.jpg","ipekyapi.com"),("Usta · balkon","/tmp/usta.jpg","ustacambalkon.com")]
for i,(cap,path,src) in enumerate(pics):
    x=.7+i*4.08; image_cover(s,path,x,2.55,3.72,2.7)
    textbox(s,cap,x,5.38,3.7,.3,12,INK,True); textbox(s,src,x,5.78,3.7,.23,8,MUTED)
textbox(s,"Ortak desen: geniş açı yaşam sahnesi + sistemin çalışmasını gösteren yakın plan + montaj/öncesi-sonrası kanıtı.",.72,6.32,11.7,.45,12,TEAL,True)
footer(s,5,"Görseller yalnız analiz referansıdır")

# 06 category visuals
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Kategori farkı","Sineklik ve premium dış mekân iki ayrı dil istiyor")
image_cover(s,"/tmp/emsa.png",.7,2.55,5.65,2.7); image_cover(s,"/tmp/winter.png",6.75,2.55,5.85,2.7)
textbox(s,"SİNEKLİK",.72,5.42,2,.3,11,TEAL,True); textbox(s,"Ürün tipi, kullanım biçimi, evcil hayvan dayanımı, kolay temizlik ve ölçü netliği.",.72,5.82,5.55,.68,12,INK)
textbox(s,"PERGOLA / KIŞ BAHÇESİ",6.77,5.42,3,.3,11,GOLD,True); textbox(s,"Mimari atmosfer, dört mevsim yaşam, ışık, manzara ve yatırım değeri.",6.77,5.82,5.55,.68,12,INK)
footer(s,6,"emsasineklik.com · wintergarden.com.tr")

# 07 web patterns
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Web yapısı","Kazandıran sayfa iskeleti")
items=[
 ("01","Sorunu göster","Yağmur, rüzgâr, böcek, kullanılmayan balkon veya güvenlik ihtiyacı."),
 ("02","Sistemi anlaşılır anlat","Teknik adı ikinci sıraya al; “yana katlanır”, “yukarı sarılır”, “dikey motorlu” de."),
 ("03","Kanıtla","Gerçek proje, ilçe, tarih, malzeme, süre, önce/sonra ve müşteri yorumu."),
 ("04","Riski azalt","Garanti, keşif, ölçü doğrulama, montaj süreci, bakım ve ödeme açıklaması."),
 ("05","Tek adım iste","Hizmet + ilçe + fotoğraf; sonra telefon. Mevcut formunuz buna çok yakın."),
]
for i,(no,h,b) in enumerate(items):
    y=2.5+i*.82; textbox(s,no,.72,y,.5,.3,10,TEAL,True); textbox(s,h,1.35,y,2.55,.33,13,INK,True); textbox(s,b,3.95,y,8.2,.45,11,MUTED)
footer(s,7)

# 08 content patterns
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Sosyal içerik","Rakibin erişimini aşacak 6 format")
formats=[
 ("10–15 sn mekanizma","Tek hareketi yakın planda göster: aç, kapat, kilitle, temizle."),
 ("Önce / sonra","Aynı kamera açısından boş alan → tamamlanmış yaşam alanı."),
 ("Usta anlatıyor","Gerçek montajcı 20 saniyede bir kritik ayrıntıyı anlatsın."),
 ("Yanlış seçim testi","Kapıya pencere modeli neden olmaz? Kedi olan evde hangi tül?"),
 ("İlçe proje günlüğü","“Beylikdüzü’nde 1 günde plise kapı uygulaması” gibi doğrulanabilir yerellik."),
 ("Müşteri sonrası","30 gün sonra kullanım, temizlik ve memnuniyet görüntüsü."),
]
for i,(h,b) in enumerate(formats):
    col=i%2; row=i//2; card(s,.7+col*6.05,2.5+row*1.32,5.65,1.05,h,b,TEAL if col==0 else GOLD)
footer(s,8,"YouTube doğrulaması: ürün-demonstrasyon formatı öne çıkıyor")

# 09 YouTube evidence
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"YouTube sinyali","Kısa ürün gösterimi, marka filminden daha verimli")
stats=[("235.854","Uzaktan kumandalı cam balkon","Albert Genau · 87 sn"),("147.113","İspanyolet kilit mekanizması","Albert Genau · 10 sn"),("10.573","Isıcamlı sistem ses yalıtımı","Albert Genau · 16 sn"),("3.949","Plisseri tanıtım filmi","Plisseri · 122 sn")]
for i,(num,h,src) in enumerate(stats):
    x=.7+(i%2)*6.05; y=2.55+(i//2)*1.65
    textbox(s,num,x,y,2.05,.58,27,TEAL,True,"Aptos Display"); textbox(s,h,x+2.25,y+.05,3.2,.36,13,INK,True); textbox(s,src,x+2.25,y+.48,3.2,.25,9,MUTED)
textbox(s,"Çıkarım: erişim için “tek fayda + tek hareket + yakın plan” kurgusu; satış için videodan doğrudan ilgili hizmet/ilçe teklif sayfasına geçiş.",.72,6.15,11.5,.55,13,INK,True)
footer(s,9,"Kaynak: yt-dlp ile 10.08.2026 tarihli arama görünümü")

# 10 opportunity
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Uma Yapı fırsatı","Rakipleri kopyalamadan daha iyi bir sistem")
bullet_list(s,[
 "Her hizmette 8–12 özgün gerçek proje; geniş açı, detay, kullanım ve montaj kareleri.",
 "Teknik ürün adı + halkın anlayacağı kullanım adı birlikte: “Plise · yana katlanır sineklik”.",
 "Her proje sayfasında ilçe, ihtiyaç, çözüm, kullanılan sistem, süre ve doğrulanmış sonuç.",
 "Hizmet seçilince doğru örneklerin açıldığı mevcut carousel’i bütün 8 hizmette eksiksiz doldurma.",
 "Teklif formunu reklam kaynağı, kreatif, hizmet, ilçe ve satış sonucu ile uçtan uca bağlama.",
 "Yalnız gerçek proje geldikçe hizmet–ilçe sayfalarını index’e açma; kopya yerel sayfa üretmeme.",
],.78,2.55,11.6,3.8,16)
footer(s,10)

# 11 shoot list
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Çekim standardı","Her işten minimum 12 içerik çıkar")
shoot=["Dış cephe geniş açı","İç mekân geniş açı","Açık / kapalı iki durum","Profil ve ray detayı","Kilit / kol kullanımı","Su tahliye detayı","Ustanın montaj anı","Önce / sonra aynı açı","Dikey 9:16 mekanizma","Müşteri kullanım anı","Renk / malzeme yakın plan","Proje tabelası: ilçe + sistem"]
for i,item in enumerate(shoot):
    col=i%3; row=i//3; x=.72+col*4.02; y=2.55+row*.83
    textbox(s,f"{i+1:02d}",x,y,.48,.3,10,TEAL,True); textbox(s,item,x+.58,y,3.15,.42,12,INK,True)
textbox(s,"Çıktılar: web hero 16:9 · carousel 4:3 · Instagram 4:5 · Reels 9:16 · detay kırpımı 1:1",.75,6.25,11.6,.45,12,GOLD,True)
footer(s,11)

# 12 roadmap
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"90 günlük plan","Araştırmayı büyümeye çevirme")
phases=[("0–30 gün","10 gerçek proje çekimi\n8 hizmet carousel’i\nMeta/GA4 olay doğrulama"),("31–60 gün","Haftada 3 Reels\n4 proje vaka sayfası\nİlk 5 hizmet–ilçe sayfası"),("61–90 gün","Kreatif A/B testleri\nKazanan ilçe/hizmete bütçe\nCRM satış sonucu geri besleme")]
for i,(h,b) in enumerate(phases): card(s,.72+i*4.02,2.65,3.65,2.9,h,b,TEAL if i<2 else GOLD)
textbox(s,"Öncelik: daha fazla sayfa değil, daha fazla doğrulanmış proje kanıtı + daha iyi dönüşüm ölçümü.",.75,6.05,11.7,.5,14,INK,True)
footer(s,12)

# 13 sources
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title(s,"Kaynakça","Doğrulanan ana kaynaklar", "Bağlantılar 10 Ağustos 2026 tarihinde herkese açık olarak incelendi.")
sources=[
 "albertgenau.com/tr · instagram.com/albert.genau · youtube.com/@ALBERTGENAU",
 "royalglass.com.tr · instagram.com/royalglasscomtr",
 "ustacambalkon.com · instagram.com/ustacambalkon · youtube.com/@ustacambalkon",
 "ipekyapi.com · plisseri.com · emsasineklik.com",
 "riccoyapi.com · wintergarden.com.tr · ehilyapi.com",
 "Agent Reach: Exa web search, Jina Reader, yt-dlp YouTube araması",
]
bullet_list(s,sources,.78,2.65,11.5,3.45,14)
textbox(s,"Kısıt: Instagram/Facebook OpenCLI tarayıcı eklentisi bağlı olmadığından gönderi ve etkileşim sayıları doğrudan doğrulanmadı; sunumda nicel sosyal sıralama yapılmadı.",.78,6.1,11.5,.6,11,GOLD,True)
footer(s,13)

prs.save(OUT)
print(OUT)
