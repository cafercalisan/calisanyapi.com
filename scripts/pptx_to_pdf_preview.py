from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "docs" / "calisan-yapi-rakip-arastirmasi-2026-08.pptx"
TARGET = ROOT / "docs" / "calisan-yapi-rakip-arastirmasi-2026-08.pdf"
PT_PER_EMU = 72 / 914400

font_regular = Path("/System/Library/Fonts/Supplemental/Arial.ttf")
font_bold = Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf")
pdfmetrics.registerFont(TTFont("CYRegular", str(font_regular)))
pdfmetrics.registerFont(TTFont("CYBold", str(font_bold)))


def rgb(value, fallback=(0.08, 0.12, 0.13)):
    try:
        if value and value.rgb:
            raw = str(value.rgb)
            return Color(int(raw[0:2], 16) / 255, int(raw[2:4], 16) / 255, int(raw[4:6], 16) / 255)
    except Exception:
        pass
    return Color(*fallback)


def wrap(text, font, size, width):
    lines = []
    for paragraph in text.splitlines() or [""]:
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            trial = current + " " + word
            if pdfmetrics.stringWidth(trial, font, size) <= width:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


prs = Presentation(SOURCE)
page_w = prs.slide_width * PT_PER_EMU
page_h = prs.slide_height * PT_PER_EMU
pdf = canvas.Canvas(str(TARGET), pagesize=(page_w, page_h), pageCompression=1)

for slide in prs.slides:
    try:
        fill = slide.background.fill.fore_color
        pdf.setFillColor(rgb(fill, (0.94, 0.92, 0.88)))
    except Exception:
        pdf.setFillColor(Color(.94, .92, .88))
    pdf.rect(0, 0, page_w, page_h, stroke=0, fill=1)

    for shape in slide.shapes:
        x = shape.left * PT_PER_EMU
        y = page_h - (shape.top + shape.height) * PT_PER_EMU
        w = shape.width * PT_PER_EMU
        h = shape.height * PT_PER_EMU

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                with Image.open(BytesIO(blob)) as im:
                    converted = BytesIO()
                    im.convert("RGB").save(converted, "JPEG", quality=88)
                    converted.seek(0)
                    pdf.drawImage(ImageReader(converted), x, y, w, h, preserveAspectRatio=False, mask="auto")
            except Exception:
                pass
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if shape.fill.type:
                    pdf.setFillColor(rgb(shape.fill.fore_color, (1, 1, 1)))
                    pdf.setStrokeColor(rgb(shape.line.color, (.8, .8, .76)))
                    pdf.rect(x, y, w, h, stroke=1, fill=1)
            except Exception:
                pass

        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue

        tf = shape.text_frame
        cursor_y = y + h - 3
        for paragraph in tf.paragraphs:
            text = paragraph.text
            if not text:
                cursor_y -= 4
                continue
            first = paragraph.runs[0] if paragraph.runs else None
            size = (first.font.size.pt if first and first.font.size else 12)
            bold = bool(first and first.font.bold)
            font = "CYBold" if bold else "CYRegular"
            color = rgb(first.font.color if first else None)
            line_height = size * 1.18
            for line in wrap(text, font, size, max(8, w - 5)):
                cursor_y -= line_height
                if cursor_y < y:
                    break
                pdf.setFont(font, size)
                pdf.setFillColor(color)
                pdf.drawString(x + 2, cursor_y, line)
            cursor_y -= size * .18

    pdf.showPage()

pdf.save()
print(TARGET)
